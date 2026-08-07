"""Generate PROJECT_SPEC.pdf and Kutaar_PPT.pptx for hackathon submission."""
import os
import unicodedata
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "submission")


def ascii_safe(text):
    """Strip or replace non-latin-1 characters for built-in PDF fonts."""
    result = []
    for ch in text:
        try:
            ch.encode("latin-1")
            result.append(ch)
        except UnicodeEncodeError:
            name = unicodedata.name(ch, "")
            if "EM DASH" in name or "HORIZONTAL BAR" in name:
                result.append("--")
            elif "EN DASH" in name:
                result.append("-")
            elif "LEFT SINGLE QUOTATION" in name:
                result.append("'")
            elif "RIGHT SINGLE QUOTATION" in name:
                result.append("'")
            elif "LEFT DOUBLE QUOTATION" in name:
                result.append('"')
            elif "RIGHT DOUBLE QUOTATION" in name:
                result.append('"')
            elif "BULLET" in name:
                result.append("*")
            elif "MULTIPLICATION" in name:
                result.append("x")
            elif "RIGHTWARDS ARROW" in name:
                result.append("->")
            elif "ELLIPSIS" in name:
                result.append("...")
            else:
                result.append("")  # drop emoji and other chars
    return "".join(result)


# ─── PDF Generation ───────────────────────────────────────────────────────────

class SpecPDF(FPDF):
    def header(self):
        self.set_font("Helvetica", "B", 10)
        self.set_text_color(100, 100, 100)
        self.cell(0, 6, ascii_safe("Kutaar -- Project Specification Document"), align="C")
        self.ln(8)

    def footer(self):
        self.set_y(-15)
        self.set_font("Helvetica", "I", 8)
        self.set_text_color(128, 128, 128)
        self.cell(0, 10, f"Page {self.page_no()}/{{nb}}", align="C")

    def section_title(self, title):
        self.set_font("Helvetica", "B", 14)
        self.set_text_color(220, 80, 0)
        self.cell(0, 8, ascii_safe(title))
        self.ln(10)

    def sub_title(self, title):
        self.set_font("Helvetica", "B", 11)
        self.set_text_color(50, 50, 50)
        self.cell(0, 7, ascii_safe(title))
        self.ln(8)

    def body_text(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 40, 40)
        self.multi_cell(0, 5.5, ascii_safe(text))
        self.ln(2)

    def bullet(self, text):
        self.set_font("Helvetica", "", 10)
        self.set_text_color(40, 40, 40)
        self.cell(5)
        self.cell(3, 5.5, "-")
        self.multi_cell(0, 5.5, ascii_safe(text))
        self.ln(1)

    def simple_table(self, headers, rows, col_widths=None):
        if col_widths is None:
            col_widths = [190 / len(headers)] * len(headers)
        self.set_font("Helvetica", "B", 9)
        self.set_fill_color(220, 80, 0)
        self.set_text_color(255, 255, 255)
        for i, h in enumerate(headers):
            self.cell(col_widths[i], 7, f" {ascii_safe(h)}", border=1, fill=True)
        self.ln()
        self.set_font("Helvetica", "", 9)
        self.set_text_color(40, 40, 40)
        for row in rows:
            for i, cell in enumerate(row):
                self.cell(col_widths[i], 6, f" {ascii_safe(cell)}", border=1)
            self.ln()
        self.ln(4)


def generate_pdf():
    pdf = SpecPDF()
    pdf.alias_nb_pages()
    pdf.set_auto_page_break(auto=True, margin=20)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_text_color(220, 80, 0)
    pdf.cell(0, 12, "Kutaar", align="C")
    pdf.ln(10)
    pdf.set_font("Helvetica", "", 12)
    pdf.set_text_color(80, 80, 80)
    pdf.cell(0, 7, "Project Specification Document", align="C")
    pdf.ln(12)
    pdf.set_font("Helvetica", "", 10)
    pdf.cell(0, 6, "AMD AI DevMaster Hackathon -- Track 2", align="C")
    pdf.ln(6)
    pdf.cell(0, 6, "Team: Script-Kitty01  |  Date: August 6, 2026", align="C")
    pdf.ln(12)

    # 1. Application Scenarios
    pdf.section_title("1. Application Scenarios")
    pdf.body_text(
        "Kutaar is a conversational multi-agent AI engineering assistant that brings the power "
        "of a full engineering review team to any developer's local machine. It runs entirely on "
        "AMD Radeon GPUs via ROCm, ensuring complete data privacy -- no code ever leaves the user's machine."
    )

    pdf.sub_title("Target Users")
    pdf.simple_table(
        ["User", "Scenario"],
        [
            ["Solo Developers", "Get instant code review feedback without waiting for teammates"],
            ["Small Teams", "Automated first-pass review before human PR review"],
            ["Security-Conscious Orgs", "Private vulnerability scanning without sending code to cloud APIs"],
            ["Open Source Maintainers", "Triage incoming PRs with automated quality checks"],
            ["Students & Learners", "Learn best practices through AI-guided code analysis"],
        ],
        [55, 135],
    )

    pdf.sub_title("Key Use Cases")
    for uc in [
        "Security Audit -- Find vulnerabilities: hardcoded secrets, SQL injection, command injection, weak hashing",
        "Performance Review -- Identify O(n^2) loops, memory-heavy patterns, blocking I/O, unnecessary deep copies",
        "Architecture Assessment -- Evaluate SOLID principles, coupling/cohesion, god classes, circular dependencies",
        "DevOps Readiness -- Check Dockerfile quality, hardcoded configs, localhost references, missing health checks",
        "Full Codebase Review -- All 4 agents run, Consensus agent synthesizes findings with quality score",
    ]:
        pdf.bullet(uc)

    # 2. Agent Architecture
    pdf.section_title("2. Agent Architecture")
    pdf.body_text(
        "Kutaar uses a 6-agent collaborative system orchestrated by LangGraph, with all inference "
        "running locally on AMD Radeon GPU via ROCm."
    )

    pdf.sub_title("Agent Roles")
    pdf.simple_table(
        ["Agent", "Role", "Tools", "Output"],
        [
            ["Planner", "Orchestrates analysis, decomposes queries", "--", "Structured task plan (JSON)"],
            ["Security", "Finds vulnerabilities (OWASP, CWE, secrets)", "Bandit, Semgrep", "Severity-scored findings"],
            ["Performance", "Spots bottlenecks & anti-patterns", "Code Search, Pattern Match", "Performance findings"],
            ["Architecture", "Evaluates design & modularity", "Git Analyzer, Code Search", "Architecture findings"],
            ["DevOps", "Checks deployment readiness", "Dockerfile Validator", "DevOps findings"],
            ["Consensus", "Cross-review debate & final verdict", "--", "Unified report + quality score"],
        ],
        [30, 45, 55, 60],
    )

    pdf.sub_title("Workflow")
    pdf.body_text(
        "User Query -> Planner -> RAG Retrieval -> [Security, Perf, Arch, DevOps] in parallel "
        "-> Consensus (debate + synthesis) -> Formatted Response -> User"
    )

    # 3. Core Capabilities
    pdf.section_title("3. Core Capabilities")
    pdf.sub_title("3.1 Multi-Agent Collaboration")
    pdf.bullet("6 specialized agents with distinct system prompts and expertise domains")
    pdf.bullet("Cross-review debate: Consensus agent identifies conflicting findings and facilitates resolution")
    pdf.bullet("Confidence scoring: Each finding includes a 0.0-1.0 confidence score, adjusted during debate")

    pdf.sub_title("3.2 Retrieval-Augmented Generation (RAG)")
    pdf.bullet("ChromaDB vector store with persistent local storage")
    pdf.bullet("sentence-transformers embeddings computed on ROCm GPU")
    pdf.bullet("Code files chunked by function/class boundaries for semantic retrieval")

    pdf.sub_title("3.3 Tool Invocation")
    pdf.bullet("Bandit -- Python static security analysis")
    pdf.bullet("Semgrep -- Multi-language pattern-based scanning")
    pdf.bullet("Git Analyzer -- Churn hotspots, commit history analysis")
    pdf.bullet("Dockerfile Validator -- Best practices checking")
    pdf.bullet("Code Search -- Regex-based pattern matching across codebase")

    pdf.sub_title("3.4 Conversation Memory")
    pdf.bullet("LangGraph checkpointing with MemorySaver")
    pdf.bullet("Full conversation history preserved across turns")
    pdf.bullet("Multi-turn drill-down: 'Tell me more about that SQL injection'")

    pdf.sub_title("3.5 Privacy-First Design")
    pdf.bullet("100% local inference -- no data sent to external APIs")
    pdf.bullet("All models run on AMD Radeon GPU via ROCm")
    pdf.bullet("No telemetry, no cloud dependency")

    # 4. Model & Deployment
    pdf.section_title("4. Model Introduction & Local Deployment Plan")

    pdf.sub_title("4.1 LLM: Llama 3.2 3B Instruct (Q4_K_M Quantized)")
    pdf.simple_table(
        ["Property", "Value"],
        [
            ["Base Model", "Meta Llama 3.2 3B Instruct"],
            ["Quantization", "Q4_K_M (4-bit with medium quality)"],
            ["Format", "GGUF"],
            ["Size on Disk", "~2.0 GB"],
            ["VRAM Usage", "~2.6 GB"],
            ["Context Window", "2048 tokens"],
            ["Batch Size", "512 tokens"],
            ["Inference Engine", "llama-cpp-python with HIP/ROCm backend"],
        ],
        [55, 135],
    )

    pdf.sub_title("4.2 Embedding Model")
    pdf.simple_table(
        ["Property", "Value"],
        [
            ["Model", "all-MiniLM-L6-v2"],
            ["Dimension", "384"],
            ["GPU Backend", "ROCm via PyTorch"],
        ],
        [55, 135],
    )

    pdf.sub_title("4.3 Quick Start")
    pdf.set_font("Courier", "", 9)
    pdf.set_text_color(40, 40, 40)
    for line in [
        "git clone https://github.com/Script-Kitty01/devMaster_amd.git",
        "cd devMaster_amd",
        "pip install -e .",
        "# Place Llama-3.2-3B-Instruct-Q4_K_M.gguf in models/",
        "python -m src.main",
    ]:
        pdf.cell(5)
        pdf.cell(0, 5, line)
        pdf.ln(5)
    pdf.ln(4)

    # 5. GPU Optimization
    pdf.section_title("5. AMD Radeon GPU Optimization")

    pdf.sub_title("5.1 HIP/ROCm Build")
    pdf.set_font("Courier", "", 8)
    for line in [
        "cmake .. -DGGML_HIP=ON -DAMDGPU_TARGETS=gfx1100",
        "  -DCMAKE_BUILD_TYPE=Release -DBUILD_SHARED_LIBS=ON",
        "  -DGGML_HIP_GRAPHS=ON -DGGML_HIP_MMQ_MFMA=ON",
    ]:
        pdf.cell(5)
        pdf.cell(0, 4.5, line)
        pdf.ln(4.5)
    pdf.ln(4)

    pdf.sub_title("5.2 Performance Benchmarks")
    pdf.simple_table(
        ["Metric", "Value"],
        [
            ["Model Load Time", "575 ms"],
            ["Prompt Eval", "57.51 ms/token (17.39 tok/s)"],
            ["Token Generation", "8.05 ms/token (124.25 tok/s)"],
            ["VRAM Usage", "2.57 GB / 51 GB"],
            ["ROCm Compute Buffer", "256.5 MiB"],
        ],
        [80, 110],
    )

    pdf.sub_title("5.3 CPU vs GPU Comparison")
    pdf.simple_table(
        ["Metric", "CPU Only", "GPU (ROCm/HIP)", "Speedup"],
        [
            ["Token Generation", "~8 tok/s", "124 tok/s", "15.5x"],
            ["Model Load", "~2.5s", "0.58s", "4.3x"],
            ["Embedding Gen", "~50 docs/s", "~400 docs/s", "8x"],
        ],
        [50, 45, 50, 45],
    )

    # 6. Technical Stack
    pdf.section_title("6. Technical Stack")
    pdf.simple_table(
        ["Layer", "Technology"],
        [
            ["GPU Runtime", "AMD ROCm 7.2.1, HIP 7.2.53211"],
            ["Inference Engine", "llama-cpp-python 0.3.34 (custom HIP build)"],
            ["LLM", "Llama 3.2 3B Instruct Q4_K_M GGUF"],
            ["Embeddings", "sentence-transformers (all-MiniLM-L6-v2) on ROCm PyTorch"],
            ["Agent Framework", "LangGraph (StateGraph + checkpointing)"],
            ["RAG", "ChromaDB (persistent local vector store)"],
            ["UI", "Streamlit"],
            ["Tools", "Bandit, Semgrep, GitPython, Docker"],
            ["Language", "Python 3.12"],
        ],
        [50, 140],
    )

    # 7. Innovation Highlights
    pdf.section_title("7. Innovation Highlights")
    for h in [
        "6-Agent Collaborative System -- Full engineering team with specialized roles, cross-review debate, and consensus building",
        "100% Local on AMD GPU -- Complete privacy; all inference, embeddings, and analysis run on the user's Radeon GPU",
        "Custom HIP Build -- llama.cpp compiled from source with GGML_HIP=ON, targeting gfx1100 with CUDA graphs and MFMA optimizations",
        "Structured Multi-Turn Memory -- LangGraph checkpointing preserves full conversation context, enabling deep drill-down analysis",
        "Tool-Augmented Agents -- Agents invoke real static analysis tools (Bandit, Semgrep) for evidence-based findings",
    ]:
        pdf.bullet(h)

    pdf.ln(10)
    pdf.set_font("Helvetica", "I", 9)
    pdf.set_text_color(128, 128, 128)
    pdf.cell(0, 6, "Built for AMD AI DevMaster Hackathon 2026 -- Track 2", align="C")

    path = os.path.join(OUTPUT_DIR, "PROJECT_SPEC.pdf")
    pdf.output(path)
    print(f"PDF saved: {path}")


# ─── PPT Generation ───────────────────────────────────────────────────────────

def add_slide(prs, layout_idx, title_text, content_lines, title_color=None):
    slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
    if title_color is None:
        title_color = RGBColor(220, 80, 0)
    title = slide.shapes.title
    title.text = title_text
    for run in title.text_frame.paragraphs[0].runs:
        run.font.color.rgb = title_color
        run.font.size = Pt(32)
        run.font.bold = True

    if content_lines and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for i, line in enumerate(content_lines):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = line
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(60, 60, 60)
            p.space_after = Pt(8)
    return slide


def add_table_slide(prs, title_text, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # blank
    left = Inches(0.5)
    top = Inches(0.3)
    width = Inches(9)
    height = Inches(0.8)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 80, 0)

    rows_count = len(rows) + 1
    cols_count = len(headers)
    table_shape = slide.shapes.add_table(
        rows_count, cols_count, Inches(0.5), Inches(1.3), Inches(9), Inches(0.4 * rows_count)
    )
    table = table_shape.table

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for par in cell.text_frame.paragraphs:
            par.font.size = Pt(14)
            par.font.bold = True
            par.font.color.rgb = RGBColor(255, 255, 255)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(220, 80, 0)

    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = val
            for par in cell.text_frame.paragraphs:
                par.font.size = Pt(13)
                par.font.color.rgb = RGBColor(50, 50, 50)
    return slide


def generate_ppt():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Kutaar"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(220, 80, 0)
    subtitle.text = (
        "Your Private AI Engineering Team\n\n"
        "AMD AI DevMaster Hackathon 2026 -- Track 2\n"
        "Team: Script-Kitty01"
    )

    # Slide 2: The Problem
    add_slide(prs, 1, "The Problem", [
        "Code review is slow, expensive, and inconsistent",
        "Cloud AI tools leak your source code",
        "Solo devs & small teams lack review bandwidth",
        "",
        "Solution: A private AI engineering team",
        "running locally on AMD Radeon",
    ])

    # Slide 3: What is Kutaar?
    add_slide(prs, 1, "What is Kutaar?", [
        "Conversational multi-agent AI assistant",
        "6 specialized agents working together",
        "100% local on AMD Radeon GPU via ROCm",
        "Upload a repo -> ask questions -> get comprehensive review",
        "",
        "No cloud. No API keys. No data leaks.",
    ])

    # Slide 4: Agent Architecture
    add_table_slide(prs, "Agent Architecture",
        ["Agent", "Role", "Tools"],
        [
            ["Planner", "Orchestrates analysis, decomposes queries", "--"],
            ["Security", "Finds vulnerabilities (OWASP, CWE, secrets)", "Bandit, Semgrep"],
            ["Performance", "Spots bottlenecks & anti-patterns", "Code Search"],
            ["Architecture", "Evaluates design & modularity", "Git Analyzer"],
            ["DevOps", "Checks deployment readiness", "Dockerfile Validator"],
            ["Consensus", "Cross-review debate & final verdict", "--"],
        ])

    # Slide 5: Core Capabilities
    add_slide(prs, 1, "Core Capabilities", [
        "Multi-turn conversation with memory (LangGraph checkpointing)",
        "RAG-powered code retrieval (ChromaDB + ROCm embeddings)",
        "Tool invocation (Bandit, Semgrep, Git, Docker)",
        "Multi-agent collaboration with cross-review debate",
        "Quality scoring (0.0-1.0) for every finding",
        "Privacy-first: 100% local, no cloud dependency",
    ])

    # Slide 6: GPU Optimization
    add_table_slide(prs, "GPU Optimization & Performance",
        ["Metric", "Value"],
        [
            ["GPU", "AMD Radeon gfx1100, 96 CUs, 51 GB VRAM"],
            ["ROCm Version", "ROCm 7.2.1 / HIP 7.2.53211"],
            ["Model", "Llama 3.2 3B Instruct Q4_K_M"],
            ["Token Generation", "124 tok/s"],
            ["Prompt Eval", "17.4 tok/s"],
            ["VRAM Usage", "2.57 GB / 51 GB"],
            ["CPU vs GPU Speedup", "15.5x faster"],
        ])

    # Slide 7: Demo Results
    add_slide(prs, 1, "Demo Results", [
        "Test repo with intentional vulnerabilities analyzed",
        "30 findings: 3 critical, 8 high, 4 medium, 15 low",
        "Overall Quality Score: 0.75",
        "",
        "Example query: 'Find security vulnerabilities'",
        "-> 10 findings, Quality Score 0.85",
        "-> Detected: hardcoded secrets, SQL injection, XSS",
    ])

    # Slide 8: Why AMD Radeon?
    add_slide(prs, 1, "Why AMD Radeon?", [
        "Complete data privacy -- no cloud, no API keys",
        "Consumer GPU accessible (4GB+ VRAM sufficient)",
        "ROCm ecosystem maturity and growing",
        "Competitive inference speed (124 tok/s on 3B model)",
        "Cost-effective vs cloud API subscriptions",
        "Open-source stack from GPU driver to inference",
    ])

    # Slide 9: Future Roadmap
    add_slide(prs, 1, "Future Roadmap", [
        "Multi-repo comparison & cross-project analysis",
        "Custom rule injection for organization-specific policies",
        "CI/CD integration (GitHub Actions, GitLab CI)",
        "Fine-tuned models specialized for code review",
        "VS Code extension for in-editor analysis",
        "Support for larger models (7B, 13B) on multi-GPU setups",
    ])

    # Slide 10: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "Thank You!"
    for run in title.text_frame.paragraphs[0].runs:
        run.font.size = Pt(48)
        run.font.bold = True
        run.font.color.rgb = RGBColor(220, 80, 0)
    subtitle.text = (
        "GitHub: Script-Kitty01/devMaster_amd\n"
        "Demo Video: https://drive.google.com/drive/folders/1PHfQT8CkQi6C3Jq6fY6C_3cnXVxyGfLl\n\n"
        "Built with love on AMD Radeon"
    )

    path = os.path.join(OUTPUT_DIR, "Kutaar_PPT.pptx")
    prs.save(path)
    print(f"PPT saved: {path}")


if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    generate_pdf()
    generate_ppt()
    print("Done! Both files generated in submission/")
